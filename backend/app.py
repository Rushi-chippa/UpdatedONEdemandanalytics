from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
import os
from werkzeug.utils import secure_filename
from analyzer import analyze_dataset
from ai_manager import get_ai_insight
from ml_engine import train_predictive_model, run_prediction
import base64
from io import BytesIO
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader
from reportlab.lib import colors
from datetime import datetime

# Ensure we setup a static path to serve the frontend folder correctly
app = Flask(__name__, static_folder="../frontend")
CORS(app)

THEME_COLORS = {
    'blue': {
        'hex': '#1e40af',
        'rgb': (30, 64, 175),
        'accent_hex': '#1e3a8a',
        'accent_rgb': (30, 58, 138)
    },
    'multicolor': {
        'hex': '#8b5cf6',
        'rgb': (139, 92, 246),
        'accent_hex': '#5b21b6',
        'accent_rgb': (91, 33, 182)
    },
    'emerald': {
        'hex': '#059669',
        'rgb': (5, 150, 105),
        'accent_hex': '#064e3b',
        'accent_rgb': (6, 78, 59)
    },
    'sunset': {
        'hex': '#ea580c',
        'rgb': (234, 88, 12),
        'accent_hex': '#7c2d12',
        'accent_rgb': (124, 45, 18)
    },
    'amethyst': {
        'hex': '#6d28d9',
        'rgb': (109, 40, 217),
        'accent_hex': '#4c1d95',
        'accent_rgb': (76, 29, 149)
    },
    'ruby': {
        'hex': '#be123c',
        'rgb': (190, 18, 60),
        'accent_hex': '#881337',
        'accent_rgb': (136, 19, 55)
    },
    'ocean': {
        'hex': '#0f766e',
        'rgb': (15, 118, 110),
        'accent_hex': '#115e59',
        'accent_rgb': (17, 94, 89)
    }
}

UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB Limit

# Serve main Dashboard
@app.route('/')
def index():
    return send_from_directory(app.static_folder, 'index.html')

# Serve static assets (js, css, etc.)
@app.route('/<path:path>')
def static_files(path):
    return send_from_directory(app.static_folder, path)

# The API Endpoint for Antigravity Analysis
@app.route('/api/analyze', methods=['POST'])
def analyze():
    # 1. Validation
    if 'files' not in request.files:
        return jsonify({"error": "No file part in request."}), 400
    
    files = request.files.getlist('files')
    if not files or files[0].filename == '':
        return jsonify({"error": "No file selected for uploading."}), 400
        
    # 2. File Processing
    allowed_extensions = ('.csv', '.xls', '.xlsx')
    filepaths = []
    
    for file in files:
        if file and file.filename.lower().endswith(allowed_extensions):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            filepaths.append(filepath)
            
    if not filepaths:
        return jsonify({"error": "Invalid file format. Please upload .csv or .xlsx files."}), 400

    # 3. Data Analytics Logic (file is PERSISTED for drill-down re-analysis)
    result = analyze_dataset(filepaths)
    
    # 4. Return Results
    if "error" in result:
        return jsonify(result), 400
        
    return jsonify(result)

# Drill-Down Filter Endpoint (re-analyzes with filters applied)
@app.route('/api/analyze_filtered', methods=['POST'])
def analyze_filtered():
    try:
        data = request.get_json()
        filename = data.get('filename')
        filters = data.get('filters', {})
        focus_metric = data.get('focus_metric')
        
        if not filename:
            return jsonify({"error": "No filename provided."}), 400
        
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        if not os.path.exists(filepath):
            return jsonify({"error": "Dataset not found. Please re-upload the file."}), 404
        
        result = analyze_dataset(filepath, filters=filters, focus_metric=focus_metric)
        
        if "error" in result:
            return jsonify(result), 400
            
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": f"Filter analysis failed: {str(e)}"}), 500


# NEW: Suggestion Endpoint for Type-Ahead Searching
@app.route('/api/suggest', methods=['POST'])
def suggest():
    try:
        import pandas as pd
        data = request.get_json()
        filename = data.get('filename')
        column = data.get('column')
        query = str(data.get('query', '')).strip().lower()
        
        if not filename or not column or not query:
            return jsonify([])

        filepath = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        if not os.path.exists(filepath):
            return jsonify([])

        # Lightweight data loading for speed
        ext = os.path.splitext(filepath)[1].lower()
        if ext == '.csv':
            df = pd.read_csv(filepath, usecols=[column], low_memory=False)
        else:
            df = pd.read_excel(filepath, usecols=[column])
        
        # Perform exact case-insensitive subset matching
        uniques = df[column].dropna().astype(str).unique()
        matches = [v for v in uniques if query in v.lower()]
        
        # Limit to top 10 most relevant matches
        return jsonify(matches[:10])

    except Exception as e:
        print(f"Suggestion Error: {e}")
        return jsonify([])

# NEW: Gemini Conversational Endpoint
@app.route('/api/chat', methods=['POST'])
def chat():
    data = request.json
    message = data.get('message', '')
    summary = data.get('data_summary', 'No data currently loaded.')
    sample_data = data.get('sample_data', [])
    
    response = get_ai_insight(message, summary, sample_data)
    
    try:
        if isinstance(response, str) and (response.strip().startswith("{") or response.strip().startswith("```json")):
            from ai_manager import clean_ai_json
            import json
            cleaned = clean_ai_json(response)
            parsed = json.loads(cleaned)
            return jsonify(parsed)
    except Exception as e:
        print(f"Failed to parse chat response as JSON: {e}")
        
    return jsonify({"type": "text", "message": response})

# NEW: Batch Chart Description Endpoint
@app.route('/api/describe_report', methods=['POST'])
def describe_report():
    data = request.json
    charts = data.get('charts', [])
    if not charts:
        return jsonify({})
    
    from ai_manager import get_report_descriptions
    descriptions = get_report_descriptions(charts)
    return jsonify(descriptions)
@app.route('/api/export_pdf', methods=['POST'])
def export_pdf():
    try:
        data = request.json
        mode = data.get('mode', 'visual')
        title = data.get('title', 'Executive Strategic Report')
        charts = data.get('charts', []) # List of {title, image(base64)}
        insights = data.get('insights', [])
        descriptions = data.get('descriptions', {}) # {chartTitle: description}
        theme = data.get('theme', 'blue')
        theme_cfg = THEME_COLORS.get(theme, THEME_COLORS['blue'])
        primary_color = colors.HexColor(theme_cfg['hex'])
        accent_color = colors.HexColor(theme_cfg['accent_hex'])

        buffer = BytesIO()
        c = canvas.Canvas(buffer, pagesize=A4)
        width, height = A4
        margin = 50
        y_pos = height - 60

        # Helper for bullet rendering
        def draw_bullets(text, x, y, max_w=500):
            lines = text.split('\n')
            temp_y = y
            c.setFont("Helvetica", 11)
            c.setFillColor(colors.black)
            for line in lines:
                if not line.strip(): continue
                clean_line = line.strip().lstrip('•').strip()
                c.drawString(x, temp_y, "•")
                words = clean_line.split()
                line_str = ""
                for w in words:
                    if c.stringWidth(line_str + " " + w) < max_w:
                        line_str += (" " if line_str else "") + w
                    else:
                        c.drawString(x + 15, temp_y, line_str)
                        temp_y -= 15
                        line_str = w
                c.drawString(x + 15, temp_y, line_str)
                temp_y -= 25
            return temp_y

        # --- DATA-FIRST ARCHITECTURE (NO EMPTY PAGES) ---
        y_pos = height - 50
        
        # Compact Brand Header
        c.setFillColor(primary_color)
        c.rect(0, height - 80, width, 80, fill=1, stroke=0)
        
        logo_path = os.path.join(app.static_folder, 'ProjectImages', 'data-analytics.png')
        if os.path.exists(logo_path):
            try:
                c.drawImage(logo_path, margin, height - 65, width=40, height=40, preserveAspectRatio=True, mask='auto')
                title_x = margin + 55
            except:
                title_x = margin
        else:
            title_x = margin

        c.setFillColor(colors.white)
        c.setFont("Helvetica-Bold", 20)
        c.drawString(title_x, height - 40, "DEMANDALYTICS STRATEGIC REPORT")
        c.setFont("Helvetica", 10)
        c.drawString(title_x, height - 60, f"Target: {title} | Mode: {mode.capitalize()} | Demandalytics Strategic Pulse Engine")
        
        y_pos = height - 110

        # High-Impact Executive Summary (Inline)
        if insights:
            c.setFont("Helvetica-Bold", 14)
            c.setFillColor(primary_color)
            c.drawString(margin, y_pos, "Executive Strategic Summary")
            y_pos -= 20
            
            c.setFont("Helvetica", 10)
            c.setFillColor(colors.black)
            for insight in insights[:4]: # Keep it compact to leave room for the first chart
                c.drawString(margin + 10, y_pos, f"• {insight}")
                y_pos -= 18
            y_pos -= 20

        # Render Charts
        if mode == 'visual':
            # === TWO CHARTS PER ROW (Side-by-Side) ===
            charts_per_page = 4
            chart_w = (width - (2 * margin) - 20) / 2
            chart_h = 220
            
            if y_pos < chart_h + 100:
                c.showPage()
                y_pos = height - 60
                
            for i, chart in enumerate(charts):
                if i > 0 and i % charts_per_page == 0:
                    c.showPage()
                    y_pos = height - 60
                
                # Calculate grid position
                col = i % 2
                row = (i // 2) % (charts_per_page // 2)
                
                curr_x = margin + (col * (chart_w + 20))
                curr_y = y_pos - ((row + 1) * (chart_h + 40))
                
                # Title
                c.setFont("Helvetica-Bold", 10)
                c.setFillColor(colors.HexColor("#1E293B"))
                c.drawString(curr_x, curr_y + chart_h + 5, chart['title'])
                
                try:
                    img_data = chart['image'].split(',')[1]
                    img_bytes = base64.b64decode(img_data)
                    img = Image.open(BytesIO(img_bytes))
                    c.drawImage(ImageReader(img), curr_x, curr_y, width=chart_w, height=chart_h, preserveAspectRatio=True)
                except: pass
        else:
            # === ONE CHART PER PAGE + DEEP AI ANALYSIS (Strategic) ===
            
            # Helper to calculate text height
            def get_bullets_height(text, max_w=500):
                lines = text.split('\n')
                needed = 0
                c.setFont("Helvetica", 11)
                for line in lines:
                    if not line.strip(): continue
                    clean_line = line.strip().lstrip('•').strip()
                    words = clean_line.split()
                    line_str = ""
                    for w in words:
                        if c.stringWidth(line_str + " " + w) < max_w:
                            line_str += (" " if line_str else "") + w
                        else:
                            needed += 15
                            line_str = w
                    needed += 25
                return needed

            for i, chart in enumerate(charts):
                desc = descriptions.get(chart['title'])
                if not desc or desc == "loading": 
                    desc = "Architecting strategic insights..."

                # Smart Pagination for the very first chart
                if i == 0:
                    is_limit = "temporarily resting" in str(desc).lower()
                    needed_space = 40 + 25 + 40 # Title + Pulse Header + Padding
                    if not (is_limit or "Synthesizing executive" in desc or "Architecting strategic" in desc):
                        needed_space += get_bullets_height(desc, max_w=width - (2 * margin))
                    else:
                        needed_space += 40
                    needed_space += 250 # Minimum acceptable height for the image
                    
                    if y_pos - needed_space < 60:
                        c.showPage()
                        y_pos = height - 60
                elif i > 0:
                    c.showPage()
                    y_pos = height - 60
                
                # 1. Chart Title (Prominent)
                c.setFont("Helvetica-Bold", 18)
                c.setFillColor(primary_color)
                c.drawString(margin, y_pos, chart['title'])
                y_pos -= 40
                
                # 2. Deep AI Analysis Content (WITH BULLET SUPPORT)
                desc = descriptions.get(chart['title'])
                if not desc or desc == "loading": 
                    desc = "Architecting strategic insights..."
                
                is_limit = "temporarily resting" in str(desc).lower()
                
                c.setFont("Helvetica-Bold", 12)
                c.setFillColor(colors.HexColor("#BC002D") if is_limit else accent_color)
                c.drawString(margin, y_pos, "DEMANDALYTICS STRATEGIC PULSE" if not is_limit else "AI Quota Notice")
                y_pos -= 25

                if is_limit or "Synthesizing executive" in desc or "Architecting strategic" in desc:
                    c.setFont("Helvetica-Oblique", 11)
                    c.setFillColor(colors.grey)
                    c.drawString(margin, y_pos, desc)
                    y_pos -= 40
                else:
                    # RENDER PROFESSIONAL BULLETS
                    y_pos = draw_bullets(desc, margin, y_pos)
                    y_pos -= 20
                
                # 3. Large Chart Image (NOW SECOND)
                try:
                    img_data = chart['image'].split(',')[1]
                    img_bytes = base64.b64decode(img_data)
                    img = Image.open(BytesIO(img_bytes))
                    
                    max_w = width - (2 * margin)
                    max_h = y_pos - 60 
                    w, h = img.size
                    
                    # Smart Scaling: If chart is extremely wide (due to adaptive scroll), 
                    # we prioritize a legible height over strict aspect ratio preservation
                    ratio = min(max_w/w, max_h/h)
                    new_w, new_h = w * ratio, h * ratio
                    
                    if new_h < 180 and max_h > 200: # Force minimum legible height if space allows
                        new_h = min(max_h - 20, 250)
                        new_w = max_w

                    c.drawImage(ImageReader(img), margin, y_pos - new_h, width=new_w, height=new_h, preserveAspectRatio=(new_h < 200))
                    y_pos -= new_h
                except: pass

        c.save()
        pdf_bytes = buffer.getvalue()
        buffer.close()

        from flask import make_response
        response = make_response(pdf_bytes)
        response.headers['Content-Type'] = 'application/pdf'
        response.headers['Content-Disposition'] = f'attachment; filename=Demandalytics_Report.pdf'
        return response

    except Exception as e:
        import traceback
        print(f"STABLE PDF EXPORT ERROR: {traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/export_pptx', methods=['POST'])
def export_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        import tempfile
        import base64
        import os
        from io import BytesIO

        data = request.json
        title_text   = data.get('title', 'Executive Strategic Report')
        charts       = data.get('charts', [])
        insights     = data.get('insights', [])
        descriptions = data.get('descriptions', {})
        mode         = data.get('mode', 'visual')
        theme        = data.get('theme', 'blue')

        theme_cfg = THEME_COLORS.get(theme, THEME_COLORS['blue'])
        r, g, b = theme_cfg['rgb']

        prs = Presentation()

        # ── Slide dimensions — read from prs to avoid any mismatch ────────────
        SLIDE_W  = prs.slide_width
        SLIDE_H  = prs.slide_height
        MARGIN   = Inches(0.30)
        HDR_H    = Inches(0.70)   # compact header band

        DARK_BLUE  = RGBColor(r, g, b)
        WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
        TEXT_DARK  = RGBColor(0x0F, 0x17, 0x2A)

        # ── Slide 1 : Title Slide ────────────────────────────────────────────
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title    = slide.shapes.title
        subtitle = slide.placeholders[1]

        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0.2), Inches(0.2),
                                        Inches(9.6), Inches(7.1))
        border.fill.background()
        border.line.color.rgb = DARK_BLUE
        border.line.width = Pt(2)

        title.text = "DEMANDALYTICS STRATEGIC REPORT"
        title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        subtitle.text = f"{title_text}\nGenerated by Demandalytics AI"

        # ── Slide 2 : Executive Summary ──────────────────────────────────────
        if insights:
            slide = prs.slides.add_slide(prs.slide_layouts[1])

            eborder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             Inches(0.2), Inches(0.2),
                                             Inches(9.6), Inches(7.1))
            eborder.fill.background()
            eborder.line.color.rgb = DARK_BLUE
            eborder.line.width = Pt(2)

            slide.shapes.title.text = "Executive Strategic Summary"
            body = slide.placeholders[1]
            body.text_frame.text = "Key Strategic AI Insights"
            body.text_frame.paragraphs[0].font.size = Pt(20) # Title text
            for insight in insights[:5]:
                p = body.text_frame.add_paragraph()
                p.text  = insight
                p.level = 1
                p.font.size = Pt(14) # Make bullet points smaller to prevent overflow

        # ── Chart Slides ─────────────────────────────────────────────────────
        for chart in charts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

            # — Dark blue header bar (compact 0.7") ─────────────────────────
            hdr_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             0, 0, SLIDE_W, HDR_H)
            hdr_bar.fill.solid()
            hdr_bar.fill.fore_color.rgb = DARK_BLUE
            hdr_bar.line.fill.background()

            # Brand label (left of header)
            brand = slide.shapes.add_textbox(Inches(0.20), Inches(0.06),
                                             Inches(2.5), HDR_H - Inches(0.06))
            brand.text_frame.paragraphs[0].text = "demandalytics"
            brand.text_frame.paragraphs[0].font.size  = Pt(9)
            brand.text_frame.paragraphs[0].font.bold  = True
            brand.text_frame.paragraphs[0].font.color.rgb = WHITE

            # Chart title (right of brand, larger)
            ttl = slide.shapes.add_textbox(Inches(2.0), Inches(0.07),
                                           SLIDE_W - Inches(2.2), HDR_H - Inches(0.07))
            ttl.text_frame.word_wrap = True
            ttl_p = ttl.text_frame.paragraphs[0]
            ttl_p.text = chart['title']
            ttl_p.font.size  = Pt(13)
            ttl_p.font.bold  = True
            ttl_p.font.color.rgb = WHITE

            # — Layout: chart area and optional bullets ──────────────────────
            desc      = descriptions.get(chart['title'], "")
            has_desc  = bool(desc and desc not in ("loading", ""))
            show_desc = (mode == 'analytical') and has_desc

            GAP       = Inches(0.08)               # small gap between header and chart
            chart_top = HDR_H + GAP

            if show_desc:
                # Chart: 56% of available height; bullets: remaining
                avail_h      = SLIDE_H - chart_top - MARGIN
                chart_height = int(avail_h * 0.56)
                desc_top     = chart_top + chart_height + GAP
                desc_height  = SLIDE_H - desc_top - Inches(0.05)
            else:
                # Chart fills entire area below header
                chart_height = SLIDE_H - chart_top - Inches(0.05)

            chart_left  = MARGIN
            chart_width = SLIDE_W - 2 * MARGIN

            # — Chart image (FULL WIDTH) ——————————————————————————————————
            has_img = False
            try:
                img_bytes = base64.b64decode(chart['image'].split(',')[1])
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                    tmp.write(img_bytes)
                    tmp_path = tmp.name

                slide.shapes.add_picture(tmp_path,
                                         chart_left, chart_top,
                                         width=chart_width, height=chart_height)
                os.remove(tmp_path)
                has_img = True
            except Exception as img_err:
                print(f"PPTX image error: {img_err}")

            # — AI bullets below chart (strategic mode) ──────────────────────
            if show_desc and has_img:
                # Light grey background strip bounded within slide
                safe_desc_h = min(desc_height, SLIDE_H - desc_top - Inches(0.03))
                if safe_desc_h > Inches(0.3):
                    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                MARGIN, desc_top,
                                                SLIDE_W - 2 * MARGIN, safe_desc_h)
                    bg.fill.solid()
                    bg.fill.fore_color.rgb = LIGHT_GRAY
                    bg.line.fill.background()

                    # Section label
                    lbl = slide.shapes.add_textbox(
                        MARGIN + Inches(0.10), desc_top + Inches(0.04),
                        SLIDE_W - 2 * MARGIN, Inches(0.20))
                    lbl.text_frame.paragraphs[0].text = "STRATEGIC AI ANALYSIS"
                    lbl.text_frame.paragraphs[0].font.size  = Pt(7)
                    lbl.text_frame.paragraphs[0].font.bold  = True
                    lbl.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE

                    # Two-column bullets, bounded strictly within slide
                    lines  = [l.strip().lstrip('•').strip() for l in desc.split('\n') if l.strip()]
                    col_w  = (SLIDE_W - 2 * MARGIN - Inches(0.20)) // 2
                    col1_x = MARGIN + Inches(0.10)
                    col2_x = col1_x + col_w + Inches(0.05)
                    y0     = desc_top + Inches(0.27)
                    row_h  = Inches(0.30)
                    max_y  = SLIDE_H - Inches(0.05)

                    for bi, line in enumerate(lines[:10]):
                        col_x = col1_x if bi % 2 == 0 else col2_x
                        row_y = y0 + (bi // 2) * row_h
                        if row_y + row_h > max_y:
                            break
                        b = slide.shapes.add_textbox(col_x, row_y, col_w, row_h)
                        b.text_frame.word_wrap = True
                        bp = b.text_frame.paragraphs[0]
                        bp.text = f"• {line}"
                        bp.font.size  = Pt(8)
                        bp.font.color.rgb = TEXT_DARK

        # ── Save ──────────────────────────────────────────────────────────────
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        from flask import make_response
        resp = make_response(buffer.getvalue())
        resp.headers['Content-Type'] = (
            'application/vnd.openxmlformats-officedocument.presentationml.presentation')
        resp.headers['Content-Disposition'] = 'attachment; filename=Demandalytics_Report.pptx'
        return resp

    except Exception as e:
        import traceback
        print(f"PPTX EXPORT ERROR: {traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/train_model', methods=['POST'])
def train_model():
    try:
        data = request.get_json()
        filename = data.get('filename')
        target = data.get('target')
        features = data.get('features', [])
        filters = data.get('filters', {})
        
        if not filename or not target or not features:
            return jsonify({"error": "Missing required parameters (filename, target, features)."}), 400
            
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], secure_filename(filename))
        if not os.path.exists(filepath):
            return jsonify({"error": "Dataset not found. Please re-upload your file."}), 404
            
        result = train_predictive_model(filepath, target, features, filters=filters)
        if "error" in result:
            return jsonify(result), 400
            
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": f"Model training request failed: {str(e)}"}), 500

@app.route('/api/predict', methods=['POST'])
def predict():
    try:
        data = request.get_json()
        filename = data.get('filename')
        model_path = data.get('model_path')
        user_inputs = data.get('inputs', {})
        
        if not filename or not model_path or not user_inputs:
            return jsonify({"error": "Missing required parameters (filename, model_path, inputs)."}), 400
            
        result = run_prediction(app.config['UPLOAD_FOLDER'], secure_filename(model_path), user_inputs)
        if "error" in result:
            return jsonify(result), 400
            
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": f"Prediction request failed: {str(e)}"}), 500

@app.route('/api/load_demo', methods=['POST'])
def load_demo():
    try:
        import pandas as pd
        import numpy as np
        import json
        
        # Generate clean B2B retail transaction demo data
        np.random.seed(42)
        dates = pd.date_range(start="2010-01-01", periods=15, freq="W")
        stores = ["Store A", "Store B", "Store C", "Store D"]
        depts = ["Electronics", "Apparel", "Home Office"]
        
        records = []
        for date in dates:
            for store in stores:
                for dept in depts:
                    base_sales = 45000 if dept == "Electronics" else (30000 if dept == "Apparel" else 20000)
                    sales = float(base_sales + np.random.randint(-8000, 15000))
                    
                    # Holiday sales spike
                    if date.month == 12:
                        sales *= 1.75
                    # Extreme Outlier Spike on Christmas week for Store A
                    if date.month == 12 and date.day == 24 and store == "Store A":
                        sales *= 2.1
                        
                    records.append({
                        "Date": date.strftime("%d-%m-%Y"),
                        "Store": store,
                        "Dept": dept,
                        "Weekly_Sales": round(sales, 2),
                        "Temperature_F": round(float(np.random.uniform(25.0, 85.0)), 1),
                        "CPI": round(float(210.5 + np.random.uniform(0.1, 8.5)), 3),
                        "Unemployment_Rate": round(float(6.8 + np.random.uniform(-0.8, 2.2)), 3)
                    })
                    
        df = pd.DataFrame(records)
        filename = "demandalytics_retail_demo.csv"
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        df.to_csv(filepath, index=False)
        
        result = analyze_dataset(filepath)
        if "error" in result:
            return jsonify(result), 400
            
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": f"Failed to generate demo data: {str(e)}"}), 500

@app.route('/api/feedback', methods=['POST'])
def save_feedback():
    try:
        import json
        data = request.get_json()
        name = data.get('name', 'Anonymous')
        email = data.get('email', 'not_provided@example.com')
        category = data.get('category', 'General feedback')
        message = data.get('message', '')
        
        if not message.strip():
            return jsonify({"error": "Message cannot be empty."}), 400
            
        feedback_file = os.path.join(app.config['UPLOAD_FOLDER'], 'feedback.json')
        
        # Load existing messages
        feedback_list = []
        if os.path.exists(feedback_file):
            try:
                with open(feedback_file, 'r', encoding='utf-8') as f:
                    feedback_list = json.load(f)
            except:
                feedback_list = []
                
        # Append new feedback
        new_feedback = {
            "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "name": name,
            "email": email,
            "category": category,
            "message": message
        }
        feedback_list.append(new_feedback)
        
        # Save back
        with open(feedback_file, 'w', encoding='utf-8') as f:
            json.dump(feedback_list, f, indent=4, ensure_ascii=False)
            
        return jsonify({"success": True, "message": "Feedback submitted successfully."})
    except Exception as e:
        return jsonify({"error": f"Failed to save feedback: {str(e)}"}), 500

if __name__ == '__main__':
    print("Starting Antigravity Server...")
    # Disabling reloader to prevent the infinite restart loop on some Windows environments
    app.run(host='0.0.0.0', port=5000, debug=True, use_reloader=False)
